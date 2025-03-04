import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def main():
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create a title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "FoodRunner"
    subtitle.text = "Delicious Food Delivered Fast"
    
    # Create content slides (simplified for working example)
    for i, title_text in enumerate([
        "Project Overview",
        "Key Features",
        "User Interface Design",
        "User Flow",
        "Technology Stack", 
        "Mobile App Integration",
        "Admin Dashboard",
        "Market Analysis",
        "Revenue Model",
        "Project Timeline",
        "Demo",
        "Our Team",
        "Thank You"
    ]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = title_text
        content.text = f"Content for {title_text}\n• Bullet point 1\n• Bullet point 2\n• Bullet point 3"
    
    # Save presentation
    output_path = "FoodRunner_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    main()